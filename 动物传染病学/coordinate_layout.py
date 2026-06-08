import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def build_smart_presentation(pptx_path, output_html, assets_dir):
    if not os.path.exists(assets_dir):
        os.makedirs(assets_dir)
    
    prs = Presentation(pptx_path)
    slides_html = ""
    
    for i, slide in enumerate(prs.slides):
        title = ""
        paragraphs = []
        images = []
        tables = []
        
        # 1. Extraction with better classification
        # We sort by top position to respect logical reading order
        shapes = sorted(slide.shapes, key=lambda s: (s.top, s.left))
        
        for shape in shapes:
            if shape.has_text_frame:
                text = "\n".join([p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()])
                if not text: continue
                
                # Heuristic: First large text block or placeholder is likely the title
                if not title and (shape.placeholder_format.type if shape.is_placeholder else 0) in [1, 3]:
                    title = text
                else:
                    paragraphs.append(text)
            
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_ext = shape.image.ext
                    img_name = f"slide_{i+1}_img_{len(images)+1}.{img_ext}"
                    with open(os.path.join(assets_dir, img_name), "wb") as f:
                        f.write(shape.image.blob)
                    images.append(img_name)
                except: pass
            
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_data = []
                for row in shape.table.rows:
                    table_data.append([cell.text_frame.text.strip() for cell in row.cells])
                tables.append(table_data)

        # 2. Slide Type Determination & HTML Building
        is_cover = i == 0 or "感谢" in title or "观看" in title
        section_attr = 'data-background-gradient="radial-gradient(circle, #232526 0%, #414345 100%)"'
        
        slide_html = f'<section {section_attr}>\n'
        
        # Title Header (Shared for non-cover slides)
        if not is_cover and title:
            slide_html += f'  <h2 class="slide-header">{title}</h2>\n'
        
        slide_html += '  <div class="slide-content-wrapper">\n'
        
        if is_cover:
            slide_html += f'    <div class="cover-box">\n'
            slide_html += f'      <h1 class="glitch-text">{title if title else "汇报展示"}</h1>\n'
            if paragraphs:
                slide_html += f'      <p class="subtitle">{" | ".join(paragraphs)}</p>\n'
            slide_html += f'    </div>\n'
        else:
            # Layout coordination based on content count
            has_img = len(images) > 0
            has_table = len(tables) > 0
            
            # Left Side: Text/Tables
            slide_html += f'    <div class="main-text {"full-width" if not has_img else "split-left"}">\n'
            for p in paragraphs:
                if len(p) > 100: # Long text gets its own block
                    slide_html += f'      <div class="text-block fragment fade-up">{p}</div>\n'
                else:
                    slide_html += f'      <li class="bullet fragment fade-in-then-semi-out">{p}</li>\n'
            
            for table in tables:
                slide_html += '      <table class="slide-table fragment zoom-in">\n'
                for row in table:
                    slide_html += '        <tr>' + "".join([f'<td>{c}</td>' for c in row]) + '</tr>\n'
                slide_html += '      </table>\n'
            slide_html += '    </div>\n'
            
            # Right Side: Images
            if has_img:
                grid_class = "single-img" if len(images) == 1 else "img-grid"
                slide_html += f'    <div class="image-area {grid_class}">\n'
                for img in images:
                    slide_html += f'      <div class="img-frame fragment scale-up"><img src="{assets_dir}/{img}"></div>\n'
                slide_html += '    </div>\n'
        
        slide_html += '  </div>\n' # End content-wrapper
        slide_html += '</section>\n'
        slides_html += slide_html

    # 3. Enhanced CSS for "Coordination"
    template = f"""
<!doctype html>
<html>
<head>
    <meta charset="utf-8">
    <title>协同排版演示文稿</title>
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reset.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.min.css">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/theme/night.min.css">
    <style>
        :root {{ --primary: #00f2fe; --gold: #ffd700; --bg-card: rgba(255,255,255,0.05); }}
        
        .reveal .slides {{ text-align: left; }}
        
        /* Layout Coordination */
        .slide-content-wrapper {{
            display: flex; gap: 30px; align-items: center; justify-content: center;
            height: 80%; width: 100%; margin-top: 20px;
        }}
        
        .main-text {{ flex: 1; display: flex; flex-direction: column; gap: 15px; }}
        .split-left {{ max-width: 60%; }}
        .full-width {{ max-width: 90%; margin: 0 auto; }}
        
        .image-area {{ flex: 1; display: flex; align-items: center; justify-content: center; gap: 10px; }}
        .img-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(180px, 1fr)); width: 100%; }}
        .single-img img {{ max-height: 500px; border-radius: 15px; }}
        
        .img-frame {{ 
            background: var(--bg-card); padding: 10px; border-radius: 10px; 
            border: 1px solid rgba(255,255,255,0.1); box-shadow: 0 10px 30px rgba(0,0,0,0.5);
        }}
        .img-frame img {{ width: 100%; height: auto; object-fit: contain; border-radius: 5px; }}

        /* Typography */
        .slide-header {{ 
            font-size: 2em !important; color: var(--primary); text-transform: none !important;
            border-bottom: 3px solid var(--primary); padding-bottom: 10px; width: fit-content;
        }}
        
        .text-block {{ 
            font-size: 0.6em; line-height: 1.6; color: #ccc; 
            background: var(--bg-card); padding: 20px; border-radius: 8px;
        }}
        
        .bullet {{ 
            font-size: 0.8em; list-style: none; position: relative; padding-left: 30px; 
            color: #fff; margin-bottom: 10px;
        }}
        .bullet::before {{ 
            content: '✦'; position: absolute; left: 0; color: var(--gold); 
        }}

        /* Table Styling */
        .slide-table {{ width: 100%; font-size: 0.5em; border-collapse: collapse; background: var(--bg-card); }}
        .slide-table td {{ border: 1px solid rgba(255,255,255,0.2); padding: 10px; }}
        .slide-table tr:nth-child(even) {{ background: rgba(255,255,255,0.02); }}

        /* Cover & Glitch */
        .cover-box {{ text-align: center; width: 100%; }}
        .glitch-text {{ font-size: 3.5em !important; font-weight: bold !important; color: #fff; text-shadow: 2px 2px var(--primary); }}
        .subtitle {{ font-size: 1.2em; color: var(--gold); margin-top: 20px; }}

        /* Utilities */
        ::-webkit-scrollbar {{ display: none; }}
    </style>
</head>
<body>
    <div class="reveal"><div class="slides">{slides_html}</div></div>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.js"></script>
    <script>
        Reveal.initialize({{ 
            width: 1280, height: 800, margin: 0.1,
            center: true, controls: true, progress: true, hash: true, transition: 'zoom',
            backgroundTransition: 'slide'
        }});
    </script>
</body>
</html>
"""
    with open(output_html, "w", encoding="utf-8") as f:
        f.write(template)

if __name__ == "__main__":
    build_smart_presentation("第一小组(1)(3).pptx", "presentation_smart.html", "pptx_assets")
